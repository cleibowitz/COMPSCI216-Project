"""
CS216 Final Project Presentation Builder
Goldman Sachs-style PowerPoint deck
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from pptx.util import Inches, Pt
import copy, os

# ── Color palette ────────────────────────────────────────────────────
NAVY    = RGBColor(0x00, 0x2B, 0x5C)
BLUE    = RGBColor(0x00, 0x5B, 0x99)
GOLD    = RGBColor(0xA0, 0x82, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
GRAY    = RGBColor(0x55, 0x55, 0x55)
LGRAY   = RGBColor(0xF2, 0xF2, 0xF2)
MGRAY   = RGBColor(0xCC, 0xCC, 0xCC)
RED     = RGBColor(0xC0, 0x00, 0x00)
GREEN   = RGBColor(0x00, 0x70, 0x2F)
DKGREEN = RGBColor(0x00, 0x50, 0x20)

# ── Slide dimensions (widescreen 13.33 × 7.5) ───────────────────────
W  = Inches(13.33)
H  = Inches(7.5)
M  = Inches(0.45)   # outer margin
HB = Inches(1.05)   # header bar height
FB = Inches(0.38)   # footer bar height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Font helpers ─────────────────────────────────────────────────────
FONT = "Calibri"

def _tf(shape):
    return shape.text_frame

def _set_para_spacing(para, before=0, after=0, line_sp=None):
    pPr = para._p.get_or_add_pPr()
    if before is not None:
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
        spcPts.set('val', str(int(before * 100)))
    if after is not None:
        spcAft = etree.SubElement(pPr, qn('a:spcAft'))
        spcPts = etree.SubElement(spcAft, qn('a:spcPts'))
        spcPts.set('val', str(int(after * 100)))
    if line_sp is not None:
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_sp * 100)))

def add_run(para, text, bold=False, italic=False, size=None, color=None, font=FONT):
    from pptx.util import Pt
    run = para.add_run()
    run.text = text
    rf = run.font
    rf.name   = font
    rf.bold   = bold
    rf.italic = italic
    if size:  rf.size  = Pt(size)
    if color: rf.color.rgb = color
    return run

def add_textbox(slide, text, x, y, w, h,
                bold=False, italic=False, size=12, color=BLACK,
                align=PP_ALIGN.LEFT, wrap=True, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    add_run(p, text, bold=bold, italic=italic, size=size, color=color, font=font)
    return tb

def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color, line=False):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    fill_shape(shape, color)
    if not line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
    return shape

# ── Slide layout helpers ─────────────────────────────────────────────

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)

def add_header(slide, title_text, subtitle=None):
    """Navy header bar + white title text + thin gold rule."""
    # Header rectangle
    bar = add_rect(slide, 0, 0, W, HB, NAVY)

    # Title text inside header
    tb = slide.shapes.add_textbox(M, Inches(0.10), W - 2*M, Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    add_run(p, title_text, bold=True, size=22, color=WHITE)

    # Optional subtitle in header
    if subtitle:
        tb2 = slide.shapes.add_textbox(M, Inches(0.68), W - 2*M, Inches(0.35))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        add_run(p2, subtitle, bold=False, italic=True, size=11, color=MGRAY)

    # Gold rule under header
    add_rect(slide, 0, HB, W, Inches(0.04), GOLD)

def add_footer(slide, slide_num=None, label="CS216 Everything Data  |  Duke University  |  Spring 2026"):
    """Thin footer with project label and slide number."""
    y_rule = H - FB - Inches(0.04)
    add_rect(slide, 0, y_rule, W, Inches(0.025), NAVY)

    # Left label
    add_textbox(slide, label,
                M, H - FB, Inches(8), FB,
                size=8, color=GRAY)

    # Right: slide number
    if slide_num:
        add_textbox(slide, str(slide_num),
                    W - Inches(1.2), H - FB, Inches(0.8), FB,
                    size=8, color=GRAY, align=PP_ALIGN.RIGHT)

def content_area():
    """Returns (x, y, w, h) for the main content zone below header."""
    x = M
    y = HB + Inches(0.18)
    w = W - 2*M
    h = H - y - FB - Inches(0.1)
    return x, y, w, h

def add_bullet_box(slide, items, x, y, w, h,
                   size=12.5, color=BLACK, bold_first=False,
                   indent=False, line_before=2, line_after=2):
    """Add a multi-bullet textbox. items = list of (text, level, bold)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, str):
            text, level, bld = item, 0, False
        else:
            text, level, bld = item[0], item[1] if len(item)>1 else 0, item[2] if len(item)>2 else False

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.level = level
        p.alignment = PP_ALIGN.LEFT
        _set_para_spacing(p, before=line_before, after=line_after)

        # Bullet marker
        if level == 0:
            marker = "▪  "
        else:
            marker = "   –  "

        run = p.add_run()
        run.text = marker + text
        run.font.name  = FONT
        run.font.size  = Pt(size)
        run.font.bold  = bld or (bold_first and first)
        run.font.color.rgb = color
    return tb

def add_label(slide, text, x, y, w=None, h=None, size=9.5, color=GRAY, bold=False, italic=True, align=PP_ALIGN.LEFT):
    w = w or Inches(6)
    h = h or Inches(0.28)
    add_textbox(slide, text, x, y, w, h, size=size, color=color,
                bold=bold, italic=italic, align=align)

# ── Table builder ────────────────────────────────────────────────────

def add_table(slide, headers, rows, x, y, w, h,
              col_widths=None, header_size=10, body_size=10,
              note=None):
    n_cols = len(headers)
    n_rows = len(rows)
    table = slide.shapes.add_table(n_rows + 1, n_cols, x, y, w, h).table

    # Column widths
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)

    def _cell(r, c, text, bg=None, txt_color=BLACK, bold=False,
               size=10, align=PP_ALIGN.CENTER):
        cell = table.cell(r, c)
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        _set_para_spacing(p, before=1, after=1)
        run = p.add_run()
        run.text = text
        run.font.name  = FONT
        run.font.bold  = bold
        run.font.size  = Pt(size)
        run.font.color.rgb = txt_color
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        return cell

    # Header row
    for c, h_text in enumerate(headers):
        _cell(0, c, h_text, bg=NAVY, txt_color=WHITE, bold=True, size=header_size)

    # Data rows
    for r, row in enumerate(rows):
        bg = LGRAY if r % 2 == 0 else WHITE
        for c, val in enumerate(row):
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            _cell(r + 1, c, str(val), bg=bg, size=body_size, align=align)

    if note:
        ny = y + h + Inches(0.04)
        add_label(slide, note, x, ny, w=w, size=8.5, color=GRAY, italic=True)

    return table

# ── Callout box ──────────────────────────────────────────────────────

def add_callout(slide, text, x, y, w, h, bg=LGRAY, border=NAVY,
                size=12, color=BLACK, bold=False):
    box = slide.shapes.add_shape(1, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = border
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_para_spacing(p, before=4, after=4, line_sp=14)
    add_run(p, text, bold=bold, size=size, color=color)

def add_stat_block(slide, x, y, w, h, stat, label, color=NAVY):
    """Large stat + small label stacked in a box."""
    box = slide.shapes.add_shape(1, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = LGRAY
    box.line.color.rgb = color
    box.line.width = Pt(2)

    tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.05),
                                   w - Inches(0.2), h - Inches(0.1))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, stat, bold=True, size=26, color=color)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    _set_para_spacing(p2, before=2)
    add_run(p2, label, bold=False, size=9, color=GRAY)

# ── Image helper ─────────────────────────────────────────────────────

def add_image(slide, path, x, y, w, h=None):
    if not os.path.exists(path):
        # Placeholder rectangle
        r = add_rect(slide, x, y, w, h or Inches(3), LGRAY)
        add_textbox(slide, f"[Figure: {os.path.basename(path)}]",
                    x + Inches(0.1), y + Inches(0.1),
                    w - Inches(0.2), h or Inches(3),
                    size=9, color=GRAY, italic=True)
        return r
    if h:
        return slide.shapes.add_picture(path, x, y, w, h)
    else:
        return slide.shapes.add_picture(path, x, y, w)

# ════════════════════════════════════════════════════════════════════
# SLIDE DEFINITIONS
# ════════════════════════════════════════════════════════════════════

FIGURES = "outputs/figures"

# ── Slide 1: Title ──────────────────────────────────────────────────
def slide_title(prs):
    s = blank_slide(prs)

    # Full navy background
    add_rect(s, 0, 0, W, H, NAVY)

    # Gold accent bar at bottom
    add_rect(s, 0, H - Inches(0.55), W, Inches(0.55), GOLD)

    # Main title
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.9), Inches(2.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_para_spacing(p, line_sp=26)
    add_run(p, "Mean Reversion in OIS-Treasury Spreads",
            bold=True, size=32, color=WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    _set_para_spacing(p2, before=4, line_sp=26)
    add_run(p2, "Evidence, Regime Dependence, and a Hidden Markov Model Trading Filter",
            bold=False, size=20, color=MGRAY)

    # Thin gold divider
    add_rect(s, Inches(0.7), Inches(3.55), Inches(4.5), Inches(0.05), GOLD)

    # Authors
    authors = [
        "Chase Leibowitz (crl65)      Dhruv Trivedi (dpt10)",
        "Jose Maldonado (jam348)   Felipe Sanchez-Noguera (fs172)",
    ]
    for i, a in enumerate(authors):
        add_textbox(s, a, Inches(0.7), Inches(3.72) + Inches(0.38)*i,
                    Inches(9), Inches(0.4),
                    size=13, color=WHITE, bold=False)

    # Course info
    add_textbox(s, "CS216 Everything Data  |  Duke University  |  April 2026",
                Inches(0.7), Inches(4.65), Inches(9), Inches(0.35),
                size=11, color=MGRAY, italic=True)

    # Video URL placeholder
    url_box = s.shapes.add_textbox(Inches(0.7), Inches(5.25), Inches(10), Inches(0.45))
    tf2 = url_box.text_frame
    p3 = tf2.paragraphs[0]
    p3.alignment = PP_ALIGN.LEFT
    add_run(p3, "Video Recording: ", bold=True, size=11, color=MGRAY)
    add_run(p3, "[INSERT VIDEO URL HERE]", bold=False, size=11, color=GOLD)

    # Bottom bar text
    add_textbox(s, "github.com/crl65/COMPSCI216-Project",
                Inches(0.7), H - Inches(0.47), Inches(8), Inches(0.38),
                size=9, color=NAVY, bold=False)

# ── Slide 2: What Is an OIS-Treasury Spread? ────────────────────────
def slide_what_is_spread(prs, n):
    s = blank_slide(prs)
    add_header(s, "The OIS-Treasury Spread: A Measure of Market Stress and Arbitrage")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    # Left column: explanation
    lw = Inches(5.8)

    # Mini headers
    def mini_header(text, y):
        add_textbox(s, text, cx, y, lw, Inches(0.3),
                    bold=True, size=12, color=NAVY)

    mini_header("What is the spread?", cy)
    bullets1 = [
        ("Treasury bonds pay a fixed yield — the government's borrowing cost", 0),
        ("SOFR OIS swaps pay the floating overnight rate compounded — near risk-free, no principal exchange", 0),
        ("Spread = OIS par rate − Treasury CMT yield  (typically negative: OIS < Treasury)", 0),
        ("Why negative? Treasuries carry a liquidity premium — investors pay up for them, bidding their yields below OIS", 0),
    ]
    yb = cy + Inches(0.30)
    tb = s.shapes.add_textbox(cx, yb, lw, Inches(1.9))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for text, lvl in bullets1:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        _set_para_spacing(p, before=2, after=2)
        add_run(p, "▪  " + text, size=11.5, color=BLACK)

    mini_header("Why should it mean-revert?", cy + Inches(2.38))
    bullets2 = [
        ("Both instruments price the same forward rate path — they cannot diverge permanently", 0),
        ("When the spread drifts far from its average, arbitrageurs enter, forcing it back", 0),
        ("Classic relative-value (RV) trade: when spread is too wide/narrow, bet on the snap-back", 0),
        ("Matters for pension funds (hedging), corporates (swap costs), and Fed market-function analysis", 0),
    ]
    yb2 = cy + Inches(2.68)
    tb2 = s.shapes.add_textbox(cx, yb2, lw, Inches(1.9))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    first = True
    for text, lvl in bullets2:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        _set_para_spacing(p, before=2, after=2)
        add_run(p, "▪  " + text, size=11.5, color=BLACK)

    # Right column: stat blocks + formula
    rx = cx + lw + Inches(0.4)
    rw = cw - lw - Inches(0.4)

    # Formula box
    add_callout(s, "Spread  =  OIS par rate  −  Treasury CMT yield",
                rx, cy, rw, Inches(0.65),
                bg=NAVY, border=NAVY, size=13, color=WHITE, bold=True)

    # Three stat blocks
    stats = [
        ("−24.7 bps", "5Y mean spread (2019–2026)"),
        ("99.9%", "Days spread is negative at 5Y/10Y"),
        ("1,890", "Daily observations in panel"),
    ]
    bh = Inches(1.18)
    by = cy + Inches(0.80)
    bw = rw
    for i, (stat, label) in enumerate(stats):
        add_stat_block(s, rx, by + i*(bh + Inches(0.12)), bw, bh,
                       stat, label, color=NAVY)

# ── Slide 3: Research Questions ──────────────────────────────────────
def slide_rqs(prs, n):
    s = blank_slide(prs)
    add_header(s, "Four Questions That Together Tell a Complete Story",
               subtitle="Introduction and Research Questions")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    rqs = [
        ("RQ1", "Mean Reversion",
         "Do OIS-Treasury spreads revert to the mean after extreme ±2σ signals?",
         "Statistical + economic significance: are moves large enough to trade?"),
        ("RQ2", "Robustness",
         "Does the result hold across all rolling windows, thresholds, and horizons?",
         "Rules out parameter selection as an explanation."),
        ("RQ3", "Regime Dependence",
         "Does reversion magnitude differ across rate-direction and volatility regimes?",
         "Identifies which macro environment makes spreads snap back hardest."),
        ("Applied Extension", "HMM + Trading",
         "Can an unsupervised Hidden Markov Model detect regimes that improve out-of-sample trading?",
         "Tests whether an ML method adds value beyond a simple rules-based strategy."),
    ]

    bw = (cw - Inches(0.3)) / 4
    colors = [NAVY, BLUE, NAVY, BLUE]

    for i, (tag, name, q, motivation) in enumerate(rqs):
        bx = cx + i * (bw + Inches(0.1))

        # Tag bar
        bar = add_rect(s, bx, cy, bw, Inches(0.42), colors[i])
        add_textbox(s, tag, bx + Inches(0.1), cy + Inches(0.06),
                    bw - Inches(0.2), Inches(0.36),
                    bold=True, size=10.5, color=WHITE, align=PP_ALIGN.CENTER)

        # Name
        add_textbox(s, name, bx + Inches(0.08), cy + Inches(0.5),
                    bw - Inches(0.16), Inches(0.38),
                    bold=True, size=12.5, color=colors[i])

        # Question
        tb = s.shapes.add_textbox(bx + Inches(0.08), cy + Inches(0.95),
                                   bw - Inches(0.16), Inches(1.55))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        _set_para_spacing(p, line_sp=14)
        add_run(p, q, size=11.5, color=BLACK)

        # Motivation divider
        add_rect(s, bx + Inches(0.08), cy + Inches(2.6),
                 bw - Inches(0.16), Inches(0.025), MGRAY)

        # Motivation
        add_textbox(s, "Why it matters: " + motivation,
                    bx + Inches(0.08), cy + Inches(2.68),
                    bw - Inches(0.16), Inches(1.6),
                    size=10, color=GRAY, italic=False)

    # Bottom callout
    add_callout(s, "These four questions are designed to escalate: prove the signal exists → prove it's robust → understand when it's strongest → extract it with ML.",
                cx, cy + ch - Inches(0.62), cw, Inches(0.55),
                bg=LGRAY, border=GOLD, size=11, color=BLACK)

# ── Slide 4: Data ────────────────────────────────────────────────────
def slide_data(prs, n):
    s = blank_slide(prs)
    add_header(s, "Panel of 1,890 Daily Observations from Two Independent Sources",
               subtitle="Data Sources and Construction")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    # Sources table
    add_textbox(s, "Data Sources", cx, cy, cw, Inches(0.28),
                bold=True, size=12, color=NAVY)

    headers = ["Source", "Series / Ticker", "Content", "Role in Analysis"]
    rows = [
        ["FRED", "DGS2, DGS5, DGS10", "U.S. Treasury CMT yields at 2Y, 5Y, 10Y maturities", "Benchmark yields"],
        ["FRED", "EFFR", "Effective Fed Funds Rate (daily)", "Rate-regime classifier"],
        ["LSEG", "USDOIS2/5/10Y=PYNY", "SOFR OIS par rates (bid/ask mid)", "Swap rates (floating leg)"],
        ["LSEG", ".MOVE", "ICE BofA MOVE Index — implied Treasury vol", "Volatility regime"],
    ]
    col_widths = [1.0, 2.2, 5.5, 2.4]
    add_table(s, headers, rows, cx, cy + Inches(0.3), cw, Inches(1.55),
              col_widths=col_widths, header_size=10, body_size=10)

    # Construction
    add_textbox(s, "Construction Pipeline", cx, cy + Inches(2.0), cw, Inches(0.28),
                bold=True, size=12, color=NAVY)

    steps = [
        "FRED holiday sentinels (\".\") → NaN; LSEG multi-row headers parsed dynamically; OIS mid = (Bid+Ask)/2",
        "Inner join by date; 5-business-day forward-fill for federal holidays (Treasury closed, OIS trades)",
        "Spread: s(t,τ) = OIS(t,τ) − DGS(t,τ)  |  Signal: 60-day rolling z-score = (spread − mean) / std",
        "2Y outliers: 23 dates with spread > 100 bps (LSEG quoting artifact, April–Oct 2022) — kept in classical tests, excluded from ML",
    ]
    ty = cy + Inches(2.30)
    tb = s.shapes.add_textbox(cx, ty, cw, Inches(1.7))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for i, step in enumerate(steps):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set_para_spacing(p, before=2, after=2)
        add_run(p, f"  {i+1}.  " + step, size=11, color=BLACK)

    # Stat boxes
    stats = [
        ("1,890", "daily obs."),
        ("2019→2026", "date range"),
        ("3", "maturities\n2Y / 5Y / 10Y"),
        ("60-day", "rolling\nz-score window"),
    ]
    bw2 = (cw - Inches(0.3)) / 4
    by2 = cy + ch - Inches(0.95)
    for i, (sv, sl) in enumerate(stats):
        bx2 = cx + i*(bw2 + Inches(0.1))
        add_stat_block(s, bx2, by2, bw2, Inches(0.88), sv, sl, color=NAVY)

# ── Slide 5: RQ1 Results ─────────────────────────────────────────────
def slide_rq1(prs, n):
    s = blank_slide(prs)
    add_header(s, "RQ1: Spreads Mean-Revert — Effect Is Large Enough to Trade",
               subtitle="Results: Mean Reversion After ±2σ Signals")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    lw = Inches(7.5)
    rw = cw - lw - Inches(0.3)
    rx = cx + lw + Inches(0.3)

    # Table
    add_textbox(s, "Table 1 — RQ1 Summary (10-day horizon, all maturities)",
                cx, cy, lw, Inches(0.28), bold=True, size=11, color=NAVY)

    headers = ["Maturity", "Signal", "n", "Mean (bps)", "Cohen's d", "t-stat", "p (BH-FDR)", "Sig"]
    rows = [
        ["2Y", "z < −2", "33", "+2.69", "+0.48", "+2.76", "0.019", "★"],
        ["2Y", "z > +2", "40", "−21.17", "−0.37", "−2.35", "0.039", "★"],
        ["5Y", "z < −2", "48", "+2.68", "+0.87", "+6.02", "<0.001", "★★★"],
        ["5Y", "z > +2", "37", "−3.85", "−0.33", "−2.00", "0.055", "—"],
        ["10Y", "z < −2", "44", "+2.94", "+0.53", "+3.53", "0.003", "★★"],
        ["10Y", "z > +2", "37", "−4.10", "−0.63", "−3.85", "0.002", "★★"],
    ]
    col_widths = [1.0, 1.0, 0.5, 1.1, 1.0, 0.95, 1.15, 0.8]
    add_table(s, headers, rows, cx, cy + Inches(0.3), lw, Inches(2.1),
              col_widths=col_widths, header_size=9.5, body_size=10,
              note="BH-FDR = Benjamini-Hochberg false discovery rate correction across all 18 cells.")

    # Key takeaways right panel
    add_textbox(s, "Key Findings", rx, cy, rw, Inches(0.28),
                bold=True, size=11, color=NAVY)

    findings = [
        "100% of 18 cells show the predicted sign",
        "12 of 18 survive FDR correction (q<0.05)",
        "5Y d=0.87 is large by conventional standards",
        "+2.68 bps exceeds 1 bp round-trip cost → tradable",
    ]
    ty = cy + Inches(0.35)
    tb = s.shapes.add_textbox(rx, ty, rw, Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for f in findings:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set_para_spacing(p, before=4, after=4)
        add_run(p, "▪  " + f, size=11, color=BLACK)

    # Callout
    add_callout(s,
                "Interpretation: When the 5Y spread dips below its 60-day average by 2 standard deviations, it rises back by an average of 2.68 bps over the next 10 days. That's a large, consistent, and economically meaningful signal — not a statistical curiosity.",
                cx, cy + ch - Inches(1.05), cw, Inches(0.95),
                bg=LGRAY, border=NAVY, size=11, color=BLACK)

# ── Slide 6: RQ1 Event Study ─────────────────────────────────────────
def slide_rq1_fig(prs, n):
    s = blank_slide(prs)
    add_header(s, "Every Panel Shows the Same V-Shape: Spreads Revert in Every Direction",
               subtitle="RQ1 Event Study — 30-Day Window Around Signal Date")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    fig_path = f"{FIGURES}/rq1_event_study.png"
    add_image(s, fig_path, cx, cy, cw - Inches(0.1), ch - Inches(0.35))

    add_label(s,
              "Each panel: average spread path in ±15-day window around the signal date, normalized to 0 at trigger. "
              "Rows = maturity; columns = direction. In every panel the spread drifts back toward zero — the expected mean-reversion shape, roughly half-complete by day 10.",
              cx, cy + ch - Inches(0.32), w=cw, size=9, italic=True, color=GRAY)

# ── Slide 7: RQ2 Robustness ──────────────────────────────────────────
def slide_rq2(prs, n):
    s = blank_slide(prs)
    add_header(s, "RQ2: Result Holds Across All 162 Parameter Combinations — This Is Structural",
               subtitle="Robustness: Rolling Window × Threshold × Horizon Grid")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    lw = Inches(4.5)
    rw = cw - lw - Inches(0.35)
    rx = cx + lw + Inches(0.35)

    # Stats table
    add_textbox(s, "RQ2 Robustness Summary (5Y spread, all horizons)",
                cx, cy, lw, Inches(0.28), bold=True, size=11, color=NAVY)

    headers = ["Metric", "Value"]
    rows = [
        ["Total parameter combinations", "162"],
        ["Cells with sign predicted by reversion", "162  (100%)"],
        ["Cells surviving BH-FDR (q<0.05)", "118  (73%)"],
        ["Range across window levels", "2.83 bps"],
        ["Range across threshold levels", "2.72 bps"],
        ["Range across horizon levels", "0.64 bps"],
    ]
    col_widths = [3.1, 1.3]
    add_table(s, headers, rows, cx, cy + Inches(0.3), lw, Inches(2.0),
              col_widths=col_widths, header_size=10, body_size=10.5)

    # Parameters tested
    add_textbox(s, "Parameter Grid", cx, cy + Inches(2.5), lw, Inches(0.28),
                bold=True, size=11, color=NAVY)
    params = [
        "Rolling window: 30, 60, 90, 120 business days",
        "Z-score threshold: |z| > 1.5, 2.0, 2.5",
        "Forward horizon: 5, 10, 20 business days",
        "Non-overlap filter: 5-day minimum gap throughout",
    ]
    tb = s.shapes.add_textbox(cx, cy + Inches(2.8), lw, Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for p_text in params:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set_para_spacing(p, before=2, after=2)
        add_run(p, "▪  " + p_text, size=11, color=BLACK)

    # Right: heatmap
    add_textbox(s, "P-Value Heatmap: All Cells Significant",
                rx, cy, rw, Inches(0.28), bold=True, size=11, color=NAVY)
    fig_path = f"{FIGURES}/rq2_heatmap.png"
    add_image(s, fig_path, rx, cy + Inches(0.3), rw, Inches(3.5))
    add_label(s, "Green = significant. Direction never flips. Magnitude barely moves.",
              rx, cy + Inches(3.85), w=rw, size=9, italic=True)

    # Callout
    add_callout(s,
                "Conclusion: Mean reversion in OIS-Treasury spreads is a structural feature of the market, not an artifact of the particular window or threshold we chose.",
                cx, cy + ch - Inches(0.62), cw, Inches(0.55),
                bg=LGRAY, border=GOLD, size=11.5, color=BLACK, bold=False)

# ── Slide 8: RQ3 ─────────────────────────────────────────────────────
def slide_rq3(prs, n):
    s = blank_slide(prs)
    add_header(s, "RQ3: Rate Regime Drives Reversion Size; Volatility Does Not",
               subtitle="Regime Dependence — Two-Way ANOVA on Rate Direction × MOVE Level")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    lw = Inches(5.5)
    rw = cw - lw - Inches(0.35)
    rx = cx + lw + Inches(0.35)

    # ANOVA table
    add_textbox(s, "Two-Way ANOVA — 5Y Spread (key rows)",
                cx, cy, lw, Inches(0.28), bold=True, size=11, color=NAVY)

    headers = ["Factor", "F-stat", "p-value", "Partial η²", "Interpretation"]
    rows = [
        ["Rate regime", "17.26", "<0.001", "0.282", "Explains 28% of variance"],
        ["Vol regime (MOVE)", "0.01", "0.930", "0.000", "Explains 0% of variance"],
        ["Rate × Vol (interaction)", "0.78", "0.383", "0.017", "No interaction effect"],
    ]
    col_widths = [1.55, 0.75, 0.85, 0.85, 1.45]
    add_table(s, headers, rows, cx, cy + Inches(0.3), lw, Inches(1.5),
              col_widths=col_widths, header_size=9.5, body_size=10.5,
              note="Partial η² = proportion of total variance explained by each factor.")

    # Rate regime split
    add_textbox(s, "Reversion by Rate Regime (5Y, z < −2, 10-day)",
                cx, cy + Inches(2.0), lw, Inches(0.28), bold=True, size=11, color=NAVY)

    headers2 = ["Rate Environment", "n", "Mean Reversion (bps)", "Cohen's d vs other"]
    rows2 = [
        ["Rising EFFR (rates climbing)", "13", "+5.40", "1.21"],
        ["Falling EFFR (rates easing)", "35", "+1.67", "—"],
    ]
    col_widths2 = [2.5, 0.45, 1.55, 1.45]
    add_table(s, headers2, rows2, cx, cy + Inches(2.3), lw, Inches(1.0),
              col_widths=col_widths2, header_size=9.5, body_size=11,
              note="Rising-rate environments produce 3× larger mean reversion.")

    # Key interpretation
    add_callout(s,
                "The answer is unambiguous: when rates are rising, the spread is under more tension — the OIS-Treasury basis is being actively repriced. That tension produces 3× larger snap-backs.\n\nVolatility level (MOVE) matters continuously (see scatter) but vanishes as a binary split. A practitioner conditioning solely on rate direction captures the regime effect cleanly.",
                cx, cy + ch - Inches(1.42), lw, Inches(1.35),
                bg=LGRAY, border=NAVY, size=10.5, color=BLACK)

    # Right: MOVE scatter
    add_textbox(s, "MOVE Level vs. Forward Reversion (continuous)",
                rx, cy, rw, Inches(0.28), bold=True, size=11, color=NAVY)
    fig_path = f"{FIGURES}/rq3_move_scatter.png"
    add_image(s, fig_path, rx, cy + Inches(0.3), rw, Inches(4.5))
    add_label(s, "OLS slope < 0, r ≈ −0.34. Higher MOVE → larger positive-z snap-backs. A simple high/low split misses this — the continuous level is needed.",
              rx, cy + Inches(4.85), w=rw, size=8.5, italic=True)

# ── Slide 9: ML Overview ─────────────────────────────────────────────
def slide_ml_overview(prs, n):
    s = blank_slide(prs)
    add_header(s, "Applied Extension: Can Machine Learning Improve on a Simple Rules-Based Entry?",
               subtitle="Two-Model Approach: Unsupervised HMM + Supervised Random Forest")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    # Three column layout
    col_w = (cw - Inches(0.4)) / 3
    labels = ["Strategy 1\nRules-Based", "Strategy 2\nRF-Filtered", "Strategy 3\nHMM-Filtered"]
    descs = [
        "Enter any time |z| > 2.\nBaseline — no ML.",
        "Enter only when Random Forest predicts reversion ≥ 3 bps.\n(Supervised, outcome labels required)",
        "Enter only when HMM assigns a 'stable' market state.\n(Unsupervised, no outcome labels)",
    ]
    colors3 = [GRAY, BLUE, NAVY]

    for i in range(3):
        bx = cx + i*(col_w + Inches(0.2))
        # Header bar
        add_rect(s, bx, cy, col_w, Inches(0.72), colors3[i])
        add_textbox(s, labels[i], bx + Inches(0.1), cy + Inches(0.1),
                    col_w - Inches(0.2), Inches(0.6),
                    bold=True, size=12, color=WHITE, align=PP_ALIGN.CENTER)
        # Description
        add_textbox(s, descs[i], bx + Inches(0.1), cy + Inches(0.8),
                    col_w - Inches(0.2), Inches(1.0),
                    size=11, color=BLACK)

    # Flow: what we're testing
    add_textbox(s, "Questions We're Answering", cx, cy + Inches(2.05), cw, Inches(0.28),
                bold=True, size=12, color=NAVY)

    qs = [
        "Does the RF's prediction of individual trade outcomes improve performance? (Supervised — needs labels)",
        "Does the HMM's detection of market regime (ignoring individual outcomes) improve performance? (Unsupervised — structure from dynamics)",
        "If both use the same regime signal, why does one approach work and the other fail?",
    ]
    tb = s.shapes.add_textbox(cx, cy + Inches(2.40), cw, Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for q in qs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set_para_spacing(p, before=4, after=4)
        add_run(p, "▪  " + q, size=11.5, color=BLACK)

    add_callout(s,
                "Central tension: the HMM regime state hurts the Random Forest as a continuous feature (−0.12 AUC) but helps the backtest as a binary gate (+0.12 Sharpe). Same signal, opposite result. Why?",
                cx, cy + ch - Inches(0.68), cw, Inches(0.6),
                bg=LGRAY, border=GOLD, size=12, color=BLACK, bold=False)

# ── Slide 10: HMM ────────────────────────────────────────────────────
def slide_hmm(prs, n):
    s = blank_slide(prs)
    add_header(s, "HMM Identifies Two Persistent, Economically Interpretable Regimes",
               subtitle="2-State Gaussian HMM — Baum-Welch Estimation, Viterbi Decoding")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    lw = Inches(5.5)
    rw = cw - lw - Inches(0.35)
    rx = cx + lw + Inches(0.35)

    # State diagnostics
    add_textbox(s, "State Diagnostics", cx, cy, lw, Inches(0.28),
                bold=True, size=11, color=NAVY)

    headers = ["", "State 0 — Stress", "State 1 — Stable"]
    rows = [
        ["Days in state", "856", "973"],
        ["% Rising EFFR", "48.1%", "18.3%"],
        ["Mean MOVE", "101.3", "77.6"],
        ["Mean 5Y Spread", "−0.269 pp", "−0.229 pp"],
        ["Avg dwell time", "~171 days", "~194 days"],
    ]
    col_widths = [1.45, 1.95, 1.95]
    add_table(s, headers, rows, cx, cy + Inches(0.3), lw, Inches(1.88),
              col_widths=col_widths, header_size=10, body_size=10.5)

    # Transition matrix
    add_textbox(s, "Transition Matrix (once in a state, very likely to stay)",
                cx, cy + Inches(2.3), lw, Inches(0.28),
                bold=True, size=11, color=NAVY)
    headers2 = ["", "→ State 0 (Stress)", "→ State 1 (Stable)"]
    rows2 = [
        ["From State 0", "0.994", "0.006"],
        ["From State 1", "0.005", "0.995"],
    ]
    col_widths2 = [1.45, 1.95, 1.95]
    add_table(s, headers2, rows2, cx, cy + Inches(2.6), lw, Inches(0.85),
              col_widths=col_widths2, header_size=10, body_size=11,
              note="States persist ~6 months — slow structural shifts a desk can actually condition on.")

    add_callout(s,
                "No look-ahead bias: refitting the HMM on pre-2024 data only produces 100% identical out-of-sample state assignments (n=591 days). The regime structure is not an in-sample artifact.",
                cx, cy + ch - Inches(0.68), lw, Inches(0.6),
                bg=LGRAY, border=GREEN, size=10.5, color=BLACK)

    # Right: HMM figure
    add_textbox(s, "Regime Decoding (Viterbi)",
                rx, cy, rw, Inches(0.28), bold=True, size=11, color=NAVY)
    fig_path = f"{FIGURES}/hmm_regimes.png"
    add_image(s, fig_path, rx, cy + Inches(0.3), rw, Inches(4.3))
    add_label(s, "Red = Stress (State 0), concentrated in high-MOVE episodes. Blue = Stable (State 1).",
              rx, cy + Inches(4.65), w=rw, size=8.5, italic=True)

# ── Slide 11: RF Failure ─────────────────────────────────────────────
def slide_rf(prs, n):
    s = blank_slide(prs)
    add_header(s, "Random Forest Fails — Not by Chance, but by Structure",
               subtitle="Supervised Classification: Why the Model Cannot Learn What the HMM Infers")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    lw = Inches(5.9)
    rw = cw - lw - Inches(0.35)
    rx = cx + lw + Inches(0.35)

    # Model comparison table
    add_textbox(s, "Out-of-Sample Performance (29 test events)",
                cx, cy, lw, Inches(0.28), bold=True, size=11, color=NAVY)

    headers = ["Model", "Accuracy", "F1", "ROC AUC"]
    rows = [
        ["Random Forest (10 features)", "72.4%", "0.833", "0.544"],
        ["Logistic Regression", "58.6%", "0.714", "0.400"],
        ["Naive (majority class)", "69.0%", "0.816", "0.500"],
    ]
    col_widths = [2.8, 1.0, 0.8, 0.9]
    add_table(s, headers, rows, cx, cy + Inches(0.3), lw, Inches(1.22),
              col_widths=col_widths, header_size=10, body_size=10.5)

    # Ablation
    add_textbox(s, "Feature Ablation: HMM State Hurts the RF",
                cx, cy + Inches(1.72), lw, Inches(0.28), bold=True, size=11, color=NAVY)

    headers2 = ["Feature Set", "Test AUC"]
    rows2 = [
        ["Full (10 features)", "0.544"],
        ["Full minus hmm_state (9 features)", "0.667"],
        ["Δ AUC from adding hmm_state", "−0.122"],
    ]
    col_widths2 = [3.8, 1.0]
    add_table(s, headers2, rows2, cx, cy + Inches(2.02), lw, Inches(1.0),
              col_widths=col_widths2, header_size=10, body_size=11)

    # Why it fails
    add_textbox(s, "Why It Fails — Three Structural Reasons", cx, cy + Inches(3.2),
                lw, Inches(0.28), bold=True, size=11, color=NAVY)

    reasons = [
        "Only 56 training events — RF has no statistical power to learn a 10-feature boundary",
        "69% majority class — RF memorizes the base rate and predicts reversion every time",
        "hmm_state is a continuous integer (0/1) — RF tries to split on it, but the regime structure only works as a binary gate, not a feature value",
    ]
    tb = s.shapes.add_textbox(cx, cy + Inches(3.5), lw, Inches(1.35))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for r in reasons:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set_para_spacing(p, before=3, after=3)
        add_run(p, "▪  " + r, size=11, color=BLACK)

    # Right: confusion matrix + key insight box
    add_textbox(s, "Confusion Matrix (OOS)", rx, cy, rw, Inches(0.28),
                bold=True, size=11, color=NAVY)
    fig_path = f"{FIGURES}/rf_confusion.png"
    add_image(s, fig_path, rx, cy + Inches(0.3), rw, Inches(2.8))

    add_callout(s,
                "RF AUC = 0.544 ≈ naive baseline 0.500.\nThe model adds no predictive value over always guessing 'yes'.",
                rx, cy + Inches(3.2), rw, Inches(0.9),
                bg=LGRAY, border=RED, size=11, color=BLACK)

    add_callout(s,
                "The HMM doesn't predict individual trade outcomes — it infers latent market dynamics directly from the spread and rate series. That's the right tool for this problem.",
                cx, cy + ch - Inches(0.65), cw, Inches(0.58),
                bg=LGRAY, border=NAVY, size=11, color=BLACK)

# ── Slide 12: Backtest Results ───────────────────────────────────────
def slide_backtest(prs, n):
    s = blank_slide(prs)
    add_header(s, "Backtest: HMM Filter Raises Sharpe by 0.12; RF Filter Does Nothing",
               subtitle="Out-of-Sample Walk-Forward Backtest — 2024-01-02 through 2026-03-24")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    lw = Inches(7.5)
    rw = cw - lw - Inches(0.3)
    rx = cx + lw + Inches(0.3)

    # Strategy comparison
    add_textbox(s, "Strategy Comparison (29 signal events, 1 bp round-trip cost)",
                cx, cy, lw, Inches(0.28), bold=True, size=11, color=NAVY)

    headers = ["Strategy", "Trades", "Win %", "Mean (bps)", "Total (bps)", "Sharpe", "Max DD", "Calmar"]
    rows = [
        ["1 — Rules-Based", "20", "55.0%", "+0.82", "+16.4", "0.39", "−9.6 bps", "0.73"],
        ["2 — RF-Filtered", "20", "55.0%", "+0.82", "+16.4", "0.39", "−9.6 bps", "0.73"],
        ["3 — HMM-Filtered", "11", "72.7%", "+1.55", "+17.1", "0.51", "−9.4 bps", "0.78"],
    ]
    col_widths = [1.9, 0.75, 0.7, 0.95, 1.0, 0.75, 0.9, 0.75]
    add_table(s, headers, rows, cx, cy + Inches(0.3), lw, Inches(1.22),
              col_widths=col_widths, header_size=9.5, body_size=10.5)

    # Per-regime breakdown
    add_textbox(s, "What the 9 Stress-State Trades Actually Did (Strategy 1, HMM state at entry)",
                cx, cy + Inches(1.72), lw, Inches(0.28), bold=True, size=11, color=NAVY)

    headers3 = ["HMM State at Entry", "Trades", "Win %", "Mean (bps)", "Total P&L (bps)"]
    rows3 = [
        ["State 1 — Stable", "11", "72.7%", "+1.55", "+17.1"],
        ["State 0 — Stress", "9", "33.3%", "−0.08", "−0.7"],
    ]
    col_widths3 = [2.4, 0.8, 0.75, 1.1, 1.4]
    add_table(s, headers3, rows3, cx, cy + Inches(2.02), lw, Inches(0.9),
              col_widths=col_widths3, header_size=10, body_size=11,
              note="All 9 stress-state trades produced all of the drawdown and zero net P&L. Strategy 3 simply skips them.")

    # Why RF is identical
    add_callout(s,
                "Why Strategy 2 = Strategy 1: The RF assigns a positive prediction to 28 of 29 test events — it passes almost everything. The filter is inoperative.",
                cx, cy + ch - Inches(0.65), lw, Inches(0.58),
                bg=LGRAY, border=RED, size=11, color=BLACK)

    # Right: stat blocks
    improvements = [
        ("0.39→0.51", "Sharpe\n(+31%)"),
        ("55%→73%", "Win rate"),
        ("55%", "Fewer\ntrades"),
    ]
    bh = Inches(1.2)
    by_start = cy + Inches(0.3)
    for i, (sv, sl) in enumerate(improvements):
        add_stat_block(s, rx, by_start + i*(bh + Inches(0.12)), rw, bh,
                       sv, sl, color=DKGREEN if i < 2 else NAVY)

    add_textbox(s, "Same entry signal.\nOnly the filter changes.",
                rx, by_start + 3*(bh + Inches(0.12)), rw, Inches(0.7),
                size=10, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

# ── Slide 13: Backtest Chart ─────────────────────────────────────────
def slide_backtest_fig(prs, n):
    s = blank_slide(prs)
    add_header(s, "HMM Filter Avoids Stress-Period Losses; RF is Indistinguishable from Rules",
               subtitle="Cumulative P&L — Three Strategies with HMM Regime Shading")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    fig_path = f"{FIGURES}/backtest_cumulative_v2.png"
    add_image(s, fig_path, cx, cy, cw - Inches(0.1), ch - Inches(0.38))

    add_label(s,
              "Background: red = HMM Stress (State 0), blue = HMM Stable (State 1). "
              "Green = HMM-filtered; orange = Rules-based; gray = RF-filtered (overlaps orange). "
              "HMM skips trades during stress periods — avoiding losses without sacrificing P&L.",
              cx, cy + ch - Inches(0.35), w=cw, size=9, italic=True, color=GRAY)

# ── Slide 14: Central Finding ────────────────────────────────────────
def slide_central_finding(prs, n):
    s = blank_slide(prs)
    add_header(s, "Representation Matters More Than Model Complexity",
               subtitle="Central Finding — The Same Signal, Two Ways to Use It")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    # Two-column contrast
    hw = (cw - Inches(0.35)) / 2
    rx = cx + hw + Inches(0.35)

    # Left: HMM as RF feature
    add_rect(s, cx, cy, hw, Inches(0.42), RED)
    add_textbox(s, "hmm_state as a Random Forest Feature",
                cx + Inches(0.1), cy + Inches(0.06), hw - Inches(0.2), Inches(0.35),
                bold=True, size=12, color=WHITE)

    items_l = [
        "Continuous integer passed into RF as one of 10 features",
        "RF tries to split on 0 vs 1 — but with 56 training samples, the split is noisy",
        "Result: AUC drops by 0.12 relative to excluding hmm_state entirely",
        "The regime signal actively hurts the classifier",
    ]
    tb = s.shapes.add_textbox(cx + Inches(0.1), cy + Inches(0.52),
                               hw - Inches(0.2), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items_l:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set_para_spacing(p, before=4, after=4)
        add_run(p, "▪  " + it, size=11.5, color=BLACK)

    add_stat_block(s, cx, cy + Inches(3.18), hw, Inches(0.95),
                   "AUC −0.12", "Adding hmm_state to RF", color=RED)

    # Right: HMM as trade gate
    add_rect(s, rx, cy, hw, Inches(0.42), DKGREEN)
    add_textbox(s, "hmm_state as a Binary Trade Gate",
                rx + Inches(0.1), cy + Inches(0.06), hw - Inches(0.2), Inches(0.35),
                bold=True, size=12, color=WHITE)

    items_r = [
        "Boolean gate: enter trade only if HMM state = Stable (State 1)",
        "No learning from trade outcomes — directly uses regime structure",
        "Correctly identifies that stress-state trades produce all drawdown, zero net P&L",
        "Result: Sharpe rises 0.39→0.51, win rate 55%→73%",
    ]
    tb2 = s.shapes.add_textbox(rx + Inches(0.1), cy + Inches(0.52),
                                hw - Inches(0.2), Inches(2.5))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    first = True
    for it in items_r:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        _set_para_spacing(p, before=4, after=4)
        add_run(p, "▪  " + it, size=11.5, color=BLACK)

    add_stat_block(s, rx, cy + Inches(3.18), hw, Inches(0.95),
                   "Sharpe +0.12", "HMM as trade gate", color=DKGREEN)

    # Central callout
    add_callout(s,
                "\"How you represent a regime signal determines whether you capture it.  Given scarce labels, imposing known market structure beats trying to learn it.\"",
                cx, cy + ch - Inches(0.95), cw, Inches(0.88),
                bg=NAVY, border=NAVY, size=14, color=WHITE, bold=True)

# ── Slide 15: Limitations & Future Work ─────────────────────────────
def slide_limitations(prs, n):
    s = blank_slide(prs)
    add_header(s, "Limitations Are Real — But the Core Signal Is Robust",
               subtitle="Limitations, Caveats, and Future Work")
    add_footer(s, n)
    cx, cy, cw, ch = content_area()

    lw = (cw - Inches(0.3)) / 2
    rx = cx + lw + Inches(0.3)

    # Limitations
    add_textbox(s, "Current Limitations", cx, cy, lw, Inches(0.28),
                bold=True, size=12, color=NAVY)

    lims = [
        ("OOS backtest window", "27 months, 29 events — directionally informative but not statistically conclusive for strategy comparisons"),
        ("ML training set", "85 total events, 56 after label construction — too small for a 10-feature classifier to learn subtle patterns"),
        ("2Y spread artifacts", "23 LSEG quoting errors in 2022 affect the 2Y maturity; excluded from ML, retained conservatively in classical tests"),
        ("Transaction costs", "Stylized flat 1 bp round-trip — real execution requires size-dependent market-impact modeling"),
    ]
    ty = cy + Inches(0.32)
    tb = s.shapes.add_textbox(cx, ty, lw, Inches(3.8))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for title_l, desc in lims:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set_para_spacing(p, before=6, after=2)
        add_run(p, title_l + ": ", bold=True, size=11.5, color=NAVY)
        add_run(p, desc, bold=False, size=11.5, color=BLACK)

    # Future work
    add_textbox(s, "Future Work", rx, cy, lw, Inches(0.28),
                bold=True, size=12, color=NAVY)

    future = [
        "Extend history to 2010 via SOFR/LIBOR-OIS splicing — 16 years vs. 7",
        "Add 1Y and 30Y maturities; test whether regime dependence scales with maturity",
        "Replace batch HMM with a real-time recursive filter for live deployment",
        "Replicate in EUR OIS (€STR) and GBP OIS (SONIA) — does the same pattern hold globally?",
        "Larger ML dataset: pooling across maturities gives ~250 events — sufficient for proper cross-validation",
        "Monitor systemic risk: at scale, RV strategies amplify dislocations during stress windows — the very periods they now skip",
    ]
    tb2 = s.shapes.add_textbox(rx, cy + Inches(0.32), lw, Inches(3.8))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    first = True
    for f in future:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        _set_para_spacing(p, before=5, after=2)
        add_run(p, "▪  " + f, size=11.5, color=BLACK)

    # Bottom summary
    add_callout(s,
                "Despite these limits: the mean-reversion signal is real (100% sign consistency, 73% FDR survival), robust (162 parameter combinations), and regime-dependent (partial η²=0.28). The HMM backtest result is suggestive, not conclusive — the signal needs more history.",
                cx, cy + ch - Inches(0.85), cw, Inches(0.78),
                bg=LGRAY, border=GOLD, size=11, color=BLACK)

# ════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════

slide_title(prs)                # 1
slide_what_is_spread(prs, 2)    # 2
slide_rqs(prs, 3)               # 3
slide_data(prs, 4)              # 4
slide_rq1(prs, 5)               # 5
slide_rq1_fig(prs, 6)           # 6
slide_rq2(prs, 7)               # 7
slide_rq3(prs, 8)               # 8
slide_ml_overview(prs, 9)       # 9
slide_hmm(prs, 10)              # 10
slide_rf(prs, 11)               # 11
slide_backtest(prs, 12)         # 12
slide_backtest_fig(prs, 13)     # 13
slide_central_finding(prs, 14)  # 14
slide_limitations(prs, 15)      # 15

out = "cs216_ois_presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({prs.slides.__len__()} slides)")
