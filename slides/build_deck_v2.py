"""
CS216 Final Project Presentation — v2
Matches example slide style: white background, Avenir Next, #1D293C dark,
#CCD3D8 section tabs, #E8ECEF alt rows, tables as primary content vehicle.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os, copy

# ── Color palette ──────────────────────────────────────────────────────────────
DARK       = RGBColor(0x1D, 0x29, 0x3C)   # #1D293C  – primary text / table headers
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
SECT_BG    = RGBColor(0xCC, 0xD3, 0xD8)   # #CCD3D8  – section label bg
ALT_ROW    = RGBColor(0xE8, 0xEC, 0xEF)   # #E8ECEF  – alternating table row
MID_GRAY   = RGBColor(0x88, 0x96, 0xA2)   # mid gray for captions
LIGHT_LINE = RGBColor(0xD0, 0xD8, 0xDE)   # light border lines
ACCENT     = RGBColor(0x2E, 0x6D, 0xA4)   # medium blue accent
GREEN_ACC  = RGBColor(0x1A, 0x6B, 0x3C)   # dark green for positive numbers
RED_ACC    = RGBColor(0xA0, 0x18, 0x18)   # dark red for negative / failure

# ── Layout constants (match example exactly) ────────────────────────────────────
W   = Inches(13.33)
H   = Inches(7.50)
LM  = Inches(0.17)   # left margin
RM  = Inches(0.17)   # right margin
CW  = W - LM - RM   # content width = 12.99"

TITLE_TOP   = Inches(0.07)
TITLE_H     = Inches(0.92)
CTX_TOP     = Inches(1.04)   # section-context bar top
CTX_H       = Inches(0.38)
CONT_TOP    = Inches(1.50)   # content area top
CONT_BOT    = Inches(6.62)   # content area bottom
CONT_H      = CONT_BOT - CONT_TOP   # ≈ 5.12"
TAB_TOP     = Inches(7.07)
TAB_H       = Inches(0.40)
PNUM_LEFT   = Inches(12.82)
PNUM_TOP    = Inches(6.70)

FONT = "Avenir Next"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


# ═══════════════════════════════════════════════════════════════════════════════
# PRIMITIVE HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def _set_para_fmt(para, before=0, after=0, line_sp=None, align=None):
    pPr = para._p.get_or_add_pPr()
    if before:
        el = etree.SubElement(pPr, qn('a:spcBef'))
        etree.SubElement(el, qn('a:spcPts')).set('val', str(int(before * 100)))
    if after:
        el = etree.SubElement(pPr, qn('a:spcAft'))
        etree.SubElement(el, qn('a:spcPts')).set('val', str(int(after * 100)))
    if line_sp:
        el = etree.SubElement(pPr, qn('a:lnSpc'))
        etree.SubElement(el, qn('a:spcPts')).set('val', str(int(line_sp * 100)))
    if align:
        pPr.set('algn', {'left':'l','center':'ctr','right':'r','justify':'just'}[align])

def run(para, text, bold=False, italic=False, size=None, color=DARK, font=FONT):
    r = para.add_run()
    r.text = text
    f = r.font
    f.name   = font
    f.bold   = bold
    f.italic = italic
    if size:  f.size = Pt(size)
    if color: f.color.rgb = color
    return r

def textbox(slide, text, x, y, w, h, bold=False, italic=False,
            size=12, color=DARK, align='left', wrap=True, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    _set_para_fmt(p, align=align)
    run(p, text, bold=bold, italic=italic, size=size, color=color, font=font)
    return tb

def rect(slide, x, y, w, h, fill, line_color=None, line_w=None):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color:
        shp.line.color.rgb = line_color
        if line_w: shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    return shp

def hline(slide, x, y, w, color=LIGHT_LINE, weight=0.75):
    ln = rect(slide, x, y, w, Inches(0.012), color)
    return ln


# ═══════════════════════════════════════════════════════════════════════════════
# CHROME (title / section bar / tab / page number)
# ═══════════════════════════════════════════════════════════════════════════════

def add_title(slide, text, size=28):
    tb = slide.shapes.add_textbox(LM, TITLE_TOP, CW, TITLE_H)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _set_para_fmt(p, align='left')
    run(p, text, bold=True, size=size, color=DARK)
    return tb

def add_section_bar(slide, section_label, subtitle=None):
    """Full-width section-context bar below the title."""
    bar = rect(slide, 0, CTX_TOP, W, CTX_H, SECT_BG)
    # Section label left
    textbox(slide, section_label,
            LM, CTX_TOP + Inches(0.06), Inches(7), CTX_H - Inches(0.06),
            bold=False, size=11, color=DARK)
    if subtitle:
        textbox(slide, subtitle,
                Inches(7), CTX_TOP + Inches(0.06), Inches(5.8), CTX_H - Inches(0.06),
                bold=False, italic=True, size=10, color=MID_GRAY, align='right')

def add_section_tab(slide, section_name, page_num):
    """Section tab at bottom-left + slide number bottom-right."""
    tab_w = Inches(max(1.8, len(section_name) * 0.115 + 0.4))
    r = rect(slide, 0, TAB_TOP, tab_w, TAB_H, SECT_BG)
    textbox(slide, section_name,
            Inches(0.15), TAB_TOP + Inches(0.10), tab_w - Inches(0.2), TAB_H,
            bold=False, size=11, color=DARK)
    # Page number
    textbox(slide, str(page_num),
            PNUM_LEFT, PNUM_TOP, Inches(0.57), Inches(0.40),
            bold=False, size=11, color=DARK, align='right')

def chrome(slide, title, section_label, page_num, subtitle=None, title_size=26):
    add_title(slide, title, size=title_size)
    add_section_bar(slide, section_label, subtitle)
    add_section_tab(slide, section_label, page_num)


# ═══════════════════════════════════════════════════════════════════════════════
# TABLE BUILDER
# ═══════════════════════════════════════════════════════════════════════════════

def gs_table(slide, headers, rows, x, y, w, h,
             col_widths=None, hdr_size=10, body_size=10.5,
             note=None, note_size=9, hdr_fill=DARK, hdr_color=WHITE,
             first_col_bold=False, no_header=False):
    """Build a GS-style table: dark header, alt-row shading, clean lines."""
    n_cols = len(headers)
    n_rows = len(rows)
    total_rows = (0 if no_header else 1) + n_rows

    tbl_shape = slide.shapes.add_table(total_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)

    def _set_cell(r, c, text, bg=None, txt_color=DARK, bold=False,
                  size=10.5, align='center', italic=False):
        cell = tbl.cell(r, c)
        cell.text = ''
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _set_para_fmt(p, before=1.5, after=1.5,
                      align='left' if (c == 0 and not no_header) else align)
        rn = p.add_run()
        rn.text = text
        rn.font.name  = FONT
        rn.font.bold  = bold
        rn.font.italic = italic
        rn.font.size  = Pt(size)
        rn.font.color.rgb = txt_color
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        return cell

    row_offset = 0
    if not no_header:
        for c, h_text in enumerate(headers):
            _set_cell(0, c, h_text, bg=hdr_fill, txt_color=hdr_color,
                      bold=True, size=hdr_size, align='center')
        row_offset = 1

    for r, row_data in enumerate(rows):
        bg = ALT_ROW if r % 2 == 0 else WHITE
        for c, val in enumerate(row_data):
            bold_cell = first_col_bold and c == 0
            _set_cell(r + row_offset, c, str(val), bg=bg, size=body_size,
                      bold=bold_cell, align='center')

    if note:
        textbox(slide, note,
                x, y + h + Inches(0.05), w, Inches(0.28),
                italic=True, size=note_size, color=MID_GRAY)
    return tbl_shape


# ═══════════════════════════════════════════════════════════════════════════════
# CONTENT HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def bullet_block(slide, items, x, y, w, h, size=11.5, color=DARK,
                 line_before=4, line_after=3, indent_size=10.5):
    """Multi-level bullet list. items = list of strings or (text, level)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, str):
            text, lvl = item, 0
        else:
            text, lvl = item[0], item[1]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set_para_fmt(p, before=line_before, after=line_after,
                      align='left')
        if lvl == 0:
            marker = "–   "
            sz = size
        else:
            marker = "      ·  "
            sz = indent_size
        run(p, marker + text, size=sz, color=color)
    return tb

def section_header(slide, text, x, y, w=None, size=11, color=DARK):
    """Small all-caps section sub-header above a table or block."""
    _w = w or (W - x - LM)
    textbox(slide, text.upper(),
            x, y, _w, Inches(0.30),
            bold=True, size=size, color=color)
    hline(slide, x, y + Inches(0.28), _w, color=DARK, weight=1)

def image_or_placeholder(slide, path, x, y, w, h=None):
    if os.path.exists(path):
        if h:
            return slide.shapes.add_picture(path, x, y, w, h)
        return slide.shapes.add_picture(path, x, y, w)
    else:
        r = rect(slide, x, y, w, h or Inches(3), ALT_ROW, LIGHT_LINE, 0.5)
        textbox(slide, f"[{os.path.basename(path)}]",
                x + Inches(0.1), y + (h or Inches(3))/2 - Inches(0.18),
                w - Inches(0.2), Inches(0.35),
                italic=True, size=9, color=MID_GRAY, align='center')
        return r

def stat_cell(slide, x, y, w, h, big_text, label, color=DARK):
    """Large number + small label in a shaded box."""
    rect(slide, x, y, w, h, ALT_ROW)
    hline(slide, x, y, w, color=color, weight=2)  # top accent line
    tb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.08),
                                   w - Inches(0.24), h - Inches(0.1))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    _set_para_fmt(p, align='center')
    run(p, big_text, bold=True, size=22, color=color)
    p2 = tf.add_paragraph()
    _set_para_fmt(p2, before=2, align='center')
    run(p2, label, bold=False, size=9.5, color=MID_GRAY)

def callout(slide, text, x, y, w, h, size=11, color=DARK, italic=False):
    """Subtle shaded callout with left accent line."""
    rect(slide, x, y, Inches(0.05), h, DARK)
    rect(slide, x + Inches(0.05), y, w - Inches(0.05), h, ALT_ROW)
    tb = slide.shapes.add_textbox(x + Inches(0.20), y + Inches(0.08),
                                   w - Inches(0.28), h - Inches(0.12))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    _set_para_fmt(p, line_sp=14, align='left')
    run(p, text, size=size, italic=italic, color=color)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

FIGURES = "outputs/figures"

# ─── Slide 1: Title ────────────────────────────────────────────────────────────
def s01_title():
    s = blank_slide()

    # Background: dark overlay on left side
    rect(s, 0, 0, W, H, DARK)
    rect(s, Inches(8.7), 0, W - Inches(8.7), H, RGBColor(0x26, 0x35, 0x4A))

    # Accent line
    rect(s, 0, Inches(2.85), Inches(8.5), Inches(0.04), SECT_BG)

    # Main title
    tb = s.shapes.add_textbox(Inches(0.55), Inches(1.20), Inches(8.0), Inches(1.55))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    _set_para_fmt(p, line_sp=34, align='left')
    run(p, "Mean Reversion in\nOIS-Treasury Spreads",
        bold=True, size=36, color=WHITE)

    # Subtitle
    tb2 = s.shapes.add_textbox(Inches(0.55), Inches(2.95), Inches(8.0), Inches(0.60))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    _set_para_fmt(p2, align='left')
    run(p2, "Evidence, Regime Dependence, and an HMM Trading Filter",
        bold=False, size=15, color=SECT_BG)

    # Authors block
    author_lines = [
        ("Chase Leibowitz (crl65)    Dhruv Trivedi (dpt10)", False),
        ("Jose Maldonado (jam348)   Felipe Sanchez-Noguera (fs172)", False),
    ]
    tb3 = s.shapes.add_textbox(Inches(0.55), Inches(3.75), Inches(8.0), Inches(0.80))
    tf3 = tb3.text_frame; tf3.word_wrap = True
    first = True
    for txt, bld in author_lines:
        p = tf3.paragraphs[0] if first else tf3.add_paragraph()
        first = False
        _set_para_fmt(p, before=3, align='left')
        run(p, txt, bold=bld, size=13, color=WHITE)

    # Course
    textbox(s, "CS216 Everything Data  ·  Duke University  ·  April 2026",
            Inches(0.55), Inches(4.72), Inches(8.0), Inches(0.35),
            italic=True, size=11, color=SECT_BG)

    # Video URL
    tb4 = s.shapes.add_textbox(Inches(0.55), Inches(5.28), Inches(8.0), Inches(0.40))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    run(p4, "Video: ", bold=True, size=11, color=SECT_BG)
    run(p4, "[INSERT VIDEO URL]", bold=False, size=11, color=WHITE)

    # Right side: key stat
    stat_lines = [
        ("1,890", "daily observations"),
        ("3", "maturities: 2Y / 5Y / 10Y"),
        ("162", "robustness cells — 100% pass"),
        ("7 years", "of daily data: 2019–2026"),
    ]
    rx = Inches(9.0)
    for i, (sv, sl) in enumerate(stat_lines):
        ty = Inches(1.5) + i * Inches(1.18)
        stat_cell(s, rx, ty, Inches(3.9), Inches(1.0), sv, sl,
                  color=SECT_BG)


# ─── Slide 2: What is the Spread? ──────────────────────────────────────────────
def s02_spread():
    s = blank_slide()
    chrome(s, "The OIS-Treasury Spread Is a Measure of Relative Value",
           "Introduction", 2,
           subtitle="What is the spread · Why it should mean-revert · Why it matters")

    # Left column (6.3")
    lw = Inches(6.3)
    ly = CONT_TOP

    section_header(s, "What Is the Spread?", LM, ly)
    bullet_block(s, [
        "OIS (Overnight Index Swap): a fixed-for-floating swap where the floating leg tracks compounded SOFR — no principal exchanged, near-zero credit risk",
        "Treasury CMT yield: the government's borrowing rate at standardized maturities (2Y, 5Y, 10Y)",
        "Spread  =  OIS par rate  −  Treasury CMT yield  →  almost always negative",
        "Why negative? Treasuries carry a superior liquidity premium — investors pay up for them, which bids their yields below OIS",
    ], LM, ly + Inches(0.35), lw, Inches(1.95), size=11.5)

    section_header(s, "Why Should It Mean-Revert?", LM, ly + Inches(2.42))
    bullet_block(s, [
        "Both instruments price the same forward short-rate path — a persistent gap is a mispricing",
        "Arbitrageurs enter when spreads drift beyond a threshold, forcing convergence",
        "Classic relative-value (RV) trade: when spread is too wide or narrow, bet on the snap-back",
    ], LM, ly + Inches(2.78), lw, Inches(1.45), size=11.5)

    section_header(s, "Why Does It Matter?", LM, ly + Inches(4.35))
    bullet_block(s, [
        "Pension funds: OIS-Treasury basis affects hedging ratios on liability-matching portfolios",
        "Corporates: swap-overlay costs are directly tied to this spread",
        "Fed: the basis is a real-time indicator of Treasury market functioning and stress",
    ], LM, ly + Inches(4.70), lw, Inches(1.38), size=11.5)

    # Right column: formula box + stats
    rx = LM + lw + Inches(0.35)
    rw = CW - lw - Inches(0.35)

    # Formula
    rect(s, rx, ly, rw, Inches(0.72), DARK)
    tb = s.shapes.add_textbox(rx + Inches(0.2), ly + Inches(0.12),
                               rw - Inches(0.4), Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; _set_para_fmt(p, align='center')
    run(p, "s", bold=True, size=14, color=WHITE)
    run(p, "(t,τ)", bold=False, size=12, color=SECT_BG)
    run(p, "  =  OIS", bold=False, size=14, color=WHITE)
    run(p, "(t,τ)", bold=False, size=12, color=SECT_BG)
    run(p, "  −  DGS", bold=False, size=14, color=WHITE)
    run(p, "(t,τ)", bold=False, size=12, color=SECT_BG)

    stats = [
        ("−24.7 bps", "5Y mean spread"),
        ("99.9%", "days spread is negative\n(5Y and 10Y)"),
        ("≥ 2σ", "entry threshold\nfor RV trade"),
        ("1 bp", "estimated round-trip\ntransaction cost"),
    ]
    for i, (sv, sl) in enumerate(stats):
        row = i // 2
        col = i % 2
        bx = rx + col * (rw / 2 + Inches(0.0))
        bw2 = rw / 2 - Inches(0.06)
        by = ly + Inches(0.85) + row * Inches(1.22)
        stat_cell(s, bx + col * Inches(0.06), by, bw2, Inches(1.10), sv, sl)

    callout(s, "Headline: OIS trades below Treasuries essentially 100% of the time — but the gap fluctuates. That fluctuation is what we trade.",
            rx, ly + Inches(3.45), rw, Inches(0.70), size=10.5)


# ─── Slide 3: Research Questions ───────────────────────────────────────────────
def s03_rqs():
    s = blank_slide()
    chrome(s, "Four Questions Designed to Escalate: Prove → Stress-Test → Explain → Exploit",
           "Introduction", 3)

    # 4-column table-style layout
    col_w = (CW - Inches(0.45)) / 4
    tags  = ["RQ1", "RQ2", "RQ3", "Applied Extension"]
    names = ["Mean Reversion", "Robustness", "Regime Dependence", "HMM + Trading"]
    questions = [
        "Do OIS-Treasury spreads mean-revert after extreme ±2σ signals, and is the move large enough to trade?",
        "Does the result hold when we change rolling window, threshold, and forward horizon — all 162 combinations?",
        "Does reversion magnitude differ across rate-direction and volatility regimes? Which factor dominates?",
        "Can an unsupervised HMM detect regimes that improve out-of-sample trading vs. a supervised Random Forest?",
    ]
    motivations = [
        "Before building a strategy, we need to know the signal is real and not a statistical artifact.",
        "A result that depends on parameter choice is fragile. We need to know it holds everywhere.",
        "Understanding when the signal is strongest determines position sizing and trade timing in practice.",
        "Tests whether ML adds value beyond a simple rules-based approach — and if so, what kind of ML.",
    ]

    for i in range(4):
        bx = LM + i * (col_w + Inches(0.15))
        by = CONT_TOP

        # Tag
        rect(s, bx, by, col_w, Inches(0.38), DARK)
        textbox(s, tags[i], bx + Inches(0.1), by + Inches(0.06),
                col_w - Inches(0.2), Inches(0.30),
                bold=True, size=13, color=WHITE, align='center')

        # Name
        textbox(s, names[i],
                bx, by + Inches(0.46), col_w, Inches(0.32),
                bold=True, size=13, color=DARK)

        # Divider
        hline(s, bx, by + Inches(0.84), col_w)

        # Question
        tb = s.shapes.add_textbox(bx, by + Inches(0.92), col_w, Inches(2.10))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        _set_para_fmt(p, line_sp=15, align='left')
        run(p, questions[i], size=11.5, color=DARK)

        # Motivation header
        textbox(s, "WHY THIS QUESTION",
                bx, by + Inches(3.10), col_w, Inches(0.24),
                bold=True, size=8.5, color=MID_GRAY)
        hline(s, bx, by + Inches(3.34), col_w, color=LIGHT_LINE)

        # Motivation text
        tb2 = s.shapes.add_textbox(bx, by + Inches(3.42), col_w, Inches(1.50))
        tf2 = tb2.text_frame; tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        _set_para_fmt(p2, line_sp=14, align='left')
        run(p2, motivations[i], size=10.5, italic=True, color=MID_GRAY)

    callout(s, "Taken together: we establish the effect, prove it is structural, identify when it is strongest, then test whether ML can systematically capture it.",
            LM, CONT_BOT - Inches(0.68), CW, Inches(0.60), size=11.5)


# ─── Slide 4: Data ─────────────────────────────────────────────────────────────
def s04_data():
    s = blank_slide()
    chrome(s, "Two Independent Sources, 1,890 Daily Observations, 2019–2026",
           "Data", 4, subtitle="Sources · Cleaning · Construction · Key decisions")

    # Sources table (full width)
    section_header(s, "Data Sources", LM, CONT_TOP)
    gs_table(s,
        ["Source", "Ticker / Series", "Content", "Role in This Analysis"],
        [
            ["FRED",  "DGS2, DGS5, DGS10",        "U.S. Treasury CMT yields at 2Y, 5Y, 10Y",     "Benchmark yields — the right-hand side of the spread"],
            ["FRED",  "EFFR",                       "Effective Fed Funds Rate (daily)",              "Rate-regime classifier: rising vs. falling"],
            ["LSEG",  "USDOIS2/5/10Y=PYNY",         "SOFR OIS par rates (bid/ask mid)",             "Swap rates — the left-hand side of the spread"],
            ["LSEG",  ".MOVE",                       "ICE BofA MOVE Index — implied Treasury vol",   "Volatility regime"],
        ],
        LM, CONT_TOP + Inches(0.32), CW, Inches(1.50),
        col_widths=[0.75, 2.0, 5.4, 4.74],
        hdr_size=10, body_size=10.5)

    # Construction (left column)
    lw = Inches(7.7)
    section_header(s, "Cleaning & Construction", LM, CONT_TOP + Inches(2.0))
    steps = [
        "FRED: sentinel \".\" on U.S. holidays → NaN before any calculation",
        "LSEG: multi-row CSV header parsed dynamically; OIS mid = (Bid + Ask) / 2",
        "Inner join by date; 5-day forward-fill for holidays where OIS trades but Treasury is closed",
        "Spread: s(t,τ) = OIS(t,τ) − DGS(t,τ)  |  Signal: 60-day rolling z-score",
        "2Y artifact: 23 dates with spread > 1 pp (LSEG quoting error, April–Oct 2022) → kept in classical tests, excluded from ML",
    ]
    bullet_block(s, steps, LM, CONT_TOP + Inches(2.34), lw, Inches(2.0), size=11.5)

    # Right column: stats
    rx = LM + lw + Inches(0.35)
    rw = CW - lw - Inches(0.35)

    section_header(s, "Panel Summary", rx, CONT_TOP + Inches(2.0))
    gs_table(s,
        ["Spread", "Mean", "Std", "% Neg"],
        [
            ["2Y",  "−9.3 bps",  "25.8 bps", "82%"],
            ["5Y",  "−24.7 bps", "10.0 bps", "≈100%"],
            ["10Y", "−33.6 bps", "13.3 bps", "≈100%"],
        ],
        rx, CONT_TOP + Inches(2.34), rw, Inches(1.30),
        col_widths=[0.62, 1.20, 1.20, 1.20],
        hdr_size=9.5, body_size=10.5)

    # Key design choices
    section_header(s, "Key Design Choices", LM, CONT_TOP + Inches(4.45))
    gs_table(s,
        ["Decision", "Value", "Rationale"],
        [
            ["Rolling window",        "60 business days", "≈3 months; balances reactivity and stability"],
            ["Signal threshold",      "|z| > 2",          "Standard RV entry trigger"],
            ["Non-overlap filter",    "5-day gap",        "Prevents double-counting the same episode"],
            ["Forward horizons",      "5, 10, 20 days",   "RQ2 robustness; main analysis uses 10-day"],
        ],
        LM, CONT_TOP + Inches(4.80), CW, Inches(1.35),
        col_widths=[2.4, 1.8, 8.69],
        hdr_size=10, body_size=10.5)


# ─── Slide 5: RQ1 Table ────────────────────────────────────────────────────────
def s05_rq1():
    s = blank_slide()
    chrome(s, "RQ1: Spreads Mean-Revert — the Effect Is Large Enough to Trade",
           "Results", 5, subtitle="Statistical + economic significance at 10-day horizon")

    lw = Inches(8.0)
    rw = CW - lw - Inches(0.35)
    rx = LM + lw + Inches(0.35)

    section_header(s, "Table 1 — RQ1 Summary (10-Day Horizon, All Maturities)", LM, CONT_TOP)
    gs_table(s,
        ["Maturity", "Signal", "n", "Mean ∆ (bps)", "Cohen's d", "t-stat", "p (BH-FDR)", "Sig"],
        [
            ["2Y",  "z < −2", "33", "+2.69",  "+0.48", "+2.76", "0.019",   "★"],
            ["2Y",  "z > +2", "40", "−21.17", "−0.37", "−2.35", "0.039",   "★"],
            ["5Y",  "z < −2", "48", "+2.68",  "+0.87", "+6.02", "<0.001",  "★★★"],
            ["5Y",  "z > +2", "37", "−3.85",  "−0.33", "−2.00", "0.055",   "—"],
            ["10Y", "z < −2", "44", "+2.94",  "+0.53", "+3.53", "0.003",   "★★"],
            ["10Y", "z > +2", "37", "−4.10",  "−0.63", "−3.85", "0.002",   "★★"],
        ],
        LM, CONT_TOP + Inches(0.34), lw, Inches(2.05),
        col_widths=[0.85, 0.90, 0.50, 1.15, 1.05, 0.90, 1.15, 0.50],
        hdr_size=9.5, body_size=10.5,
        note="BH-FDR = Benjamini-Hochberg false discovery rate correction applied across all 18 cells (3 maturities × 2 directions × 3 horizons).")

    # Key takeaways right
    section_header(s, "Key Findings", rx, CONT_TOP)
    stats_r = [
        ("100%", "of 18 cells show\npredicted sign"),
        ("12/18", "cells survive\nFDR correction"),
        ("0.87", "Cohen's d at 5Y\n(\"large\" effect)"),
    ]
    for i, (sv, sl) in enumerate(stats_r):
        stat_cell(s, rx, CONT_TOP + Inches(0.34) + i * Inches(1.18),
                  rw, Inches(1.06), sv, sl, color=ACCENT)

    # Interpretation
    section_header(s, "Interpretation", LM, CONT_TOP + Inches(2.55))
    gs_table(s,
        ["Maturity", "Best Signal", "Reversion (10d)", "Interpretation"],
        [
            ["5Y", "z < −2",  "+2.68 bps  (d=0.87)", "Exceeds 1 bp round-trip cost → economically tradable"],
            ["10Y","z > +2",  "−4.10 bps  (d=0.63)", "Large snap-back; most consistent two-sided signal"],
            ["2Y", "z > +2",  "−21.17 bps (d=0.37)", "Dominated by LSEG quoting artifacts — treat cautiously"],
        ],
        LM, CONT_TOP + Inches(2.90), CW, Inches(1.25),
        col_widths=[0.85, 1.20, 2.35, 8.49],
        hdr_size=10, body_size=10.5)

    callout(s, "The 5Y z<−2 signal produces a +2.68 bps mean reversion over 10 days (95% CI: [+1.81, +3.61]). That is larger than the 1 bp round-trip cost — the signal is real and tradable.",
            LM, CONT_BOT - Inches(0.65), CW, Inches(0.58), size=11.5)


# ─── Slide 6: RQ1 Event Study ──────────────────────────────────────────────────
def s06_event_study():
    s = blank_slide()
    chrome(s, "Every Maturity, Every Direction Shows the Same V-Shape Recovery",
           "Results", 6, subtitle="RQ1 Event Study — 30-day window centered on signal date")

    img_h = Inches(4.65)
    image_or_placeholder(s, f"{FIGURES}/rq1_event_study.png",
                         LM, CONT_TOP, CW, img_h)

    callout(s,
            "Each panel: average spread path in a ±15-day window, normalized to 0 at signal date. "
            "Rows = maturity; columns = direction (spread too low / too high). In every panel the spread drifts back toward zero — the expected mean-reversion shape, roughly half-complete by day 10.",
            LM, CONT_TOP + img_h + Inches(0.10), CW, Inches(0.65), size=11, italic=True)


# ─── Slide 7: RQ2 ──────────────────────────────────────────────────────────────
def s07_rq2():
    s = blank_slide()
    chrome(s, "RQ2: Result Holds Across All 162 Parameter Combinations — This Is Structural",
           "Results", 7,
           subtitle="Robustness: Window (30–120d) × Threshold (1.5–2.5σ) × Horizon (5–20d)")

    lw = Inches(4.8)
    rw = CW - lw - Inches(0.35)
    rx = LM + lw + Inches(0.35)

    section_header(s, "RQ2 Stability Summary", LM, CONT_TOP)
    gs_table(s,
        ["Metric", "Value"],
        [
            ["Total parameter combinations",              "162"],
            ["Cells with predicted sign",                 "162  (100%)"],
            ["Cells surviving BH-FDR (q < 0.05)",        "118  (73%)"],
            ["Mean reversion range across windows",       "2.83 bps"],
            ["Mean reversion range across thresholds",    "2.72 bps"],
            ["Mean reversion range across horizons",      "0.64 bps"],
        ],
        LM, CONT_TOP + Inches(0.34), lw, Inches(2.08),
        col_widths=[3.45, 1.25],
        hdr_size=10, body_size=11, first_col_bold=True)

    section_header(s, "Parameter Grid Tested", LM, CONT_TOP + Inches(2.55))
    gs_table(s,
        ["Parameter", "Values Tested"],
        [
            ["Rolling window",   "30, 60, 90, 120 business days"],
            ["Z-score threshold","|z| > 1.5,  2.0,  2.5"],
            ["Forward horizon",  "5, 10, 20 business days"],
        ],
        LM, CONT_TOP + Inches(2.90), lw, Inches(1.05),
        col_widths=[1.85, 2.85],
        hdr_size=10, body_size=11)

    callout(s,
            "Conclusion: the direction never flips once. The magnitude barely moves. Mean reversion in OIS-Treasury spreads is a structural property of the market, not a consequence of parameter selection.",
            LM, CONT_BOT - Inches(0.68), lw, Inches(0.60), size=11)

    # Right: heatmap
    section_header(s, "P-Value Heatmap — 5Y Spread, 10-Day Horizon", rx, CONT_TOP)
    image_or_placeholder(s, f"{FIGURES}/rq2_heatmap.png",
                         rx, CONT_TOP + Inches(0.34), rw, Inches(3.65))
    textbox(s, "Green = BH-significant. All 9 window × threshold cells point in the same direction across every horizon.",
            rx, CONT_TOP + Inches(4.05), rw, Inches(0.50),
            italic=True, size=9.5, color=MID_GRAY)


# ─── Slide 8: RQ3 ──────────────────────────────────────────────────────────────
def s08_rq3():
    s = blank_slide()
    chrome(s, "RQ3: Rate Regime Drives Reversion 3× — Volatility Contributes Nothing",
           "Results", 8,
           subtitle="Two-way ANOVA: Rate direction × MOVE level on 5Y reversion magnitude")

    lw = Inches(5.85)
    rw = CW - lw - Inches(0.35)
    rx = LM + lw + Inches(0.35)

    section_header(s, "Two-Way ANOVA — 5Y Spread", LM, CONT_TOP)
    gs_table(s,
        ["Factor", "F-stat", "p-value", "Partial η²", "Verdict"],
        [
            ["Rate regime (rising vs. falling EFFR)", "17.26", "<0.001", "0.282", "Explains 28% of variance ← dominant"],
            ["Vol regime (MOVE above/below median)",   "0.01",  "0.930",  "0.000", "Explains 0% — no effect"],
            ["Rate × Vol interaction",                  "0.78",  "0.383",  "0.017", "No interaction effect"],
        ],
        LM, CONT_TOP + Inches(0.34), lw, Inches(1.30),
        col_widths=[2.85, 0.75, 0.80, 0.80, 1.55],
        hdr_size=9.5, body_size=10.5)

    section_header(s, "Reversion Size by Rate Regime (5Y, z < −2, 10-Day)", LM, CONT_TOP + Inches(1.82))
    gs_table(s,
        ["Rate Environment", "n", "Mean Reversion", "Cohen's d"],
        [
            ["Rising EFFR (rates climbing)", "13", "+5.40 bps", "1.21 (very large)"],
            ["Falling EFFR (rates easing)",  "35", "+1.67 bps", "—"],
        ],
        LM, CONT_TOP + Inches(2.16), lw, Inches(0.90),
        col_widths=[2.40, 0.45, 1.45, 1.45],
        hdr_size=9.5, body_size=11,
        note="Rising-rate environments produce 3× larger mean reversion (p=0.001, d=1.21 — a very large effect size).")

    callout(s,
            "Why? When rates are rising, the OIS-Treasury spread is under fundamentally different tension — the basis is being actively repriced by the market. That tension produces larger, faster snap-backs.\n\nVolatility (MOVE) matters continuously (see scatter, right) but vanishes as a binary split. Conditioning on rate direction captures the regime effect cleanly.",
            LM, CONT_TOP + Inches(3.65), lw, Inches(1.10), size=11)

    # Right: MOVE scatter
    section_header(s, "MOVE Level vs. 10-Day Reversion (Continuous)", rx, CONT_TOP)
    image_or_placeholder(s, f"{FIGURES}/rq3_move_scatter.png",
                         rx, CONT_TOP + Inches(0.34), rw, Inches(3.80))
    textbox(s, "OLS slope < 0, r ≈ −0.34 (positive-z signals). Higher vol → larger snap-back. A binary split misses this; the continuous level matters.",
            rx, CONT_TOP + Inches(4.20), rw, Inches(0.50),
            italic=True, size=9.5, color=MID_GRAY)


# ─── Slide 9: ML Overview ──────────────────────────────────────────────────────
def s09_ml_overview():
    s = blank_slide()
    chrome(s, "Can ML Improve on a Simple Rules-Based Entry? Two Approaches Tested",
           "Applied Extension", 9,
           subtitle="Supervised: Random Forest classifier · Unsupervised: 2-state Gaussian HMM")

    # Three panel layout
    panel_w = (CW - Inches(0.40)) / 3

    panels = [
        ("Strategy 1", "Rules-Based Baseline",
         "Enter any time |z| > 2.\nNo ML — pure signal.",
         "Establishes what a mechanical strategy achieves with no additional information.", DARK),
        ("Strategy 2", "Random Forest Filter",
         "Enter only when RF predicts spread reverts ≥3 bps in 10 days.\nSupervised: needs outcome labels.",
         "Tests whether predicting individual trade outcomes improves performance.", ACCENT),
        ("Strategy 3", "HMM Filter",
         "Enter only in HMM 'stable' state (State 1).\nUnsupervised: no outcome labels needed.",
         "Tests whether detecting latent market regime improves performance.", GREEN_ACC),
    ]

    for i, (tag, name, desc, why, color) in enumerate(panels):
        bx = LM + i * (panel_w + Inches(0.20))
        by = CONT_TOP
        rect(s, bx, by, panel_w, Inches(0.40), color)
        textbox(s, tag, bx + Inches(0.1), by + Inches(0.06),
                panel_w - Inches(0.2), Inches(0.32),
                bold=True, size=13, color=WHITE, align='center')
        textbox(s, name, bx, by + Inches(0.50), panel_w, Inches(0.32),
                bold=True, size=12, color=DARK)
        hline(s, bx, by + Inches(0.88), panel_w)
        tb = s.shapes.add_textbox(bx, by + Inches(0.96), panel_w, Inches(1.20))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; _set_para_fmt(p, line_sp=14, align='left')
        run(p, desc, size=11.5, color=DARK)
        textbox(s, "Why: " + why, bx, by + Inches(2.28),
                panel_w, Inches(1.10), size=10.5, italic=True, color=MID_GRAY)

    section_header(s, "The Central Tension", LM, CONT_TOP + Inches(3.50))
    gs_table(s,
        ["", "As RF feature (continuous)", "As HMM gate (binary)"],
        [
            ["hmm_state signal", "AUC drops −0.12 (hurts classifier)", "Sharpe rises +0.12 (helps trading)"],
        ],
        LM, CONT_TOP + Inches(3.84), CW, Inches(0.72),
        col_widths=[2.20, 5.39, 5.40],
        hdr_size=10, body_size=12, first_col_bold=True)

    callout(s,
            "Same signal. Opposite result. The question is why — and the answer is about how the signal is represented, not how sophisticated the model is.",
            LM, CONT_BOT - Inches(0.65), CW, Inches(0.58), size=12)


# ─── Slide 10: HMM ─────────────────────────────────────────────────────────────
def s10_hmm():
    s = blank_slide()
    chrome(s, "HMM Identifies Two Persistent, Economically Meaningful Regimes",
           "Applied Extension", 10,
           subtitle="2-state Gaussian HMM · Baum-Welch estimation · Viterbi decoding")

    lw = Inches(5.85)
    rw = CW - lw - Inches(0.35)
    rx = LM + lw + Inches(0.35)

    section_header(s, "State Diagnostics", LM, CONT_TOP)
    gs_table(s,
        ["", "State 0 — Stress", "State 1 — Stable"],
        [
            ["Days in state",       "856",     "973"],
            ["% Rising EFFR",       "48.1%",   "18.3%"],
            ["Mean MOVE",           "101.3",    "77.6"],
            ["Mean 5Y spread",      "−0.269 pp","−0.229 pp"],
            ["Avg dwell time",      "~171 days","~194 days"],
        ],
        LM, CONT_TOP + Inches(0.34), lw, Inches(1.90),
        col_widths=[1.60, 2.05, 2.10],
        hdr_size=10, body_size=11, first_col_bold=True)

    section_header(s, "Transition Matrix", LM, CONT_TOP + Inches(2.38))
    gs_table(s,
        ["", "→ State 0 (Stress)", "→ State 1 (Stable)"],
        [
            ["From State 0", "0.994", "0.006"],
            ["From State 1", "0.005", "0.995"],
        ],
        LM, CONT_TOP + Inches(2.72), lw, Inches(0.88),
        col_widths=[1.60, 2.05, 2.10],
        hdr_size=9.5, body_size=11,
        note="Once in a state, probability of staying is >99%. States persist for months — slow structural shifts a desk can actually condition on.")

    callout(s,
            "No look-ahead bias: refitting HMM on pre-2024 data only → 100% identical OOS state assignments (n=591 days). The regime structure is not an in-sample artifact.",
            LM, CONT_BOT - Inches(0.65), lw, Inches(0.58), size=11)

    # Right: HMM figure
    section_header(s, "Regime Decoding (Viterbi)", rx, CONT_TOP)
    image_or_placeholder(s, f"{FIGURES}/hmm_regimes.png",
                         rx, CONT_TOP + Inches(0.34), rw, Inches(4.20))
    textbox(s, "Top: 5Y spread colored by state. Bottom: MOVE by state. State 0 (red) concentrates in high-MOVE stress episodes.",
            rx, CONT_TOP + Inches(4.60), rw, Inches(0.55),
            italic=True, size=9.5, color=MID_GRAY)


# ─── Slide 11: RF Failure ──────────────────────────────────────────────────────
def s11_rf():
    s = blank_slide()
    chrome(s, "Random Forest Fails — Not by Chance, but by Structure",
           "Applied Extension", 11,
           subtitle="56 training events · 69% majority class · hmm_state hurts AUC")

    lw = Inches(5.85)
    rw = CW - lw - Inches(0.35)
    rx = LM + lw + Inches(0.35)

    section_header(s, "OOS Performance (29 test events)", LM, CONT_TOP)
    gs_table(s,
        ["Model", "Accuracy", "F1", "ROC AUC"],
        [
            ["Random Forest (10 features)", "72.4%", "0.833", "0.544"],
            ["Logistic Regression",          "58.6%", "0.714", "0.400"],
            ["Naive (predict majority class)","69.0%", "0.816", "0.500"],
        ],
        LM, CONT_TOP + Inches(0.34), lw, Inches(1.18),
        col_widths=[2.85, 0.90, 0.72, 0.90],
        hdr_size=10, body_size=11)

    section_header(s, "Feature Ablation: Adding hmm_state Hurts RF", LM, CONT_TOP + Inches(1.65))
    gs_table(s,
        ["Feature Set", "Test AUC", "vs. Full Model"],
        [
            ["Full model (10 features)",     "0.544", "—"],
            ["Full minus hmm_state (9 feat.)","0.667", "+0.122"],
            ["Naive baseline",               "0.500", "—"],
        ],
        LM, CONT_TOP + Inches(2.00), lw, Inches(1.10),
        col_widths=[3.05, 0.88, 1.65],
        hdr_size=10, body_size=11,
        note="Removing hmm_state improves AUC by 0.12. The regime signal actively hurts the Random Forest.")

    section_header(s, "Three Structural Reasons for Failure", LM, CONT_TOP + Inches(3.25))
    bullet_block(s, [
        "Sample size: only 56 training events — no statistical power to learn a 10-feature decision boundary",
        "Class imbalance: 69% majority class → RF memorizes the base rate and predicts 'yes' on nearly every event",
        "Wrong representation: hmm_state is a continuous integer in the RF, but it only contains information as a binary gate",
    ], LM, CONT_TOP + Inches(3.60), lw, Inches(1.60), size=11.5)

    # Right: confusion matrix
    section_header(s, "Confusion Matrix (OOS)", rx, CONT_TOP)
    image_or_placeholder(s, f"{FIGURES}/rf_confusion.png",
                         rx, CONT_TOP + Inches(0.34), rw, Inches(2.50))
    textbox(s, "RF correctly predicts 'yes' on all positives (recall=1.00) but cannot identify negatives — it is not classifying, it is guessing the majority class.",
            rx, CONT_TOP + Inches(2.90), rw, Inches(0.60),
            italic=True, size=9.5, color=MID_GRAY)

    section_header(s, "Feature Importance", rx, CONT_TOP + Inches(3.60))
    image_or_placeholder(s, f"{FIGURES}/rf_importance.png",
                         rx, CONT_TOP + Inches(3.94), rw, Inches(1.55))


# ─── Slide 12: Backtest ────────────────────────────────────────────────────────
def s12_backtest():
    s = blank_slide()
    chrome(s, "Backtest: HMM Gate Raises Sharpe 31%; RF Filter Is Inoperative",
           "Applied Extension", 12,
           subtitle="Walk-forward OOS: 2024-01-02 → 2026-03-24 · 1 bp round-trip cost · 10-day hold")

    lw = Inches(8.30)
    rw = CW - lw - Inches(0.35)
    rx = LM + lw + Inches(0.35)

    section_header(s, "Strategy Comparison — 29 Signal Events", LM, CONT_TOP)
    gs_table(s,
        ["Strategy", "Trades", "Win %", "Mean (bps)", "Total (bps)", "Sharpe", "Max DD", "Calmar"],
        [
            ["1 — Rules-Based",        "20", "55.0%", "+0.82", "+16.4", "0.39", "−9.6 bps", "0.73"],
            ["2 — RF-Filtered",        "20", "55.0%", "+0.82", "+16.4", "0.39", "−9.6 bps", "0.73"],
            ["3 — HMM-Filtered (S=1)", "11", "72.7%", "+1.55", "+17.1", "0.51", "−9.4 bps", "0.78"],
        ],
        LM, CONT_TOP + Inches(0.34), lw, Inches(1.15),
        col_widths=[2.35, 0.72, 0.72, 0.90, 0.90, 0.72, 0.92, 0.75],
        hdr_size=9.5, body_size=11)

    section_header(s, "Per-Regime Breakdown (Strategy 1 trades, grouped by HMM state at entry)", LM, CONT_TOP + Inches(1.62))
    gs_table(s,
        ["HMM State at Entry", "Trades", "Win %", "Mean P&L (bps)", "Total P&L (bps)"],
        [
            ["State 1 — Stable", "11", "72.7%", "+1.55", "+17.1"],
            ["State 0 — Stress", "9",  "33.3%", "−0.08", "−0.7"],
        ],
        LM, CONT_TOP + Inches(1.96), lw, Inches(0.88),
        col_widths=[3.00, 0.72, 0.72, 1.55, 1.55],
        hdr_size=10, body_size=11,
        note="The 9 stress-state trades contributed all of the drawdown and zero net P&L. Strategy 3 simply skips them.")

    section_header(s, "Why RF = Rules (Strategy 2 = Strategy 1)", LM, CONT_TOP + Inches(3.02))
    callout(s, "RF assigns a positive prediction to 28 of 29 OOS events — it passes everything. The filter is inoperative. This is the expected outcome given the 69% training majority class.",
            LM, CONT_TOP + Inches(3.36), lw, Inches(0.62), size=11)

    callout(s, "HMM improvement: by skipping stress-state trades, Sharpe rises 0.39 → 0.51, win rate 55% → 73%, and max drawdown marginally improves — without changing a single entry signal.",
            LM, CONT_BOT - Inches(0.68), lw, Inches(0.58), size=11.5)

    # Right: key stats
    kpi = [
        ("0.39→0.51", "Sharpe (+31%)", GREEN_ACC),
        ("55%→73%",   "Win rate", GREEN_ACC),
        ("45%", "Fewer trades\n(better efficiency)", ACCENT),
    ]
    for i, (sv, sl, co) in enumerate(kpi):
        stat_cell(s, rx, CONT_TOP + Inches(0.34) + i * Inches(1.30),
                  rw, Inches(1.18), sv, sl, color=co)


# ─── Slide 13: Backtest Chart ──────────────────────────────────────────────────
def s13_backtest_chart():
    s = blank_slide()
    chrome(s, "HMM Filter Sidesteps All Stress-Period Losses; RF Tracks the Baseline",
           "Applied Extension", 13,
           subtitle="Cumulative P&L — three strategies with HMM regime shading")

    img_h = Inches(4.50)
    image_or_placeholder(s, f"{FIGURES}/backtest_cumulative_v2.png",
                         LM, CONT_TOP, CW, img_h)

    callout(s,
            "Background shading: red = HMM Stress (State 0), blue = HMM Stable (State 1). "
            "Green (HMM-filtered) avoids stress episodes. Orange (rules-based) and RF-filtered lines overlap — the RF adds nothing. "
            "HMM's selective entry produces higher Sharpe without sacrificing total P&L.",
            LM, CONT_TOP + img_h + Inches(0.12), CW, Inches(0.70), size=11, italic=True)


# ─── Slide 14: Central Finding ─────────────────────────────────────────────────
def s14_finding():
    s = blank_slide()
    chrome(s, "Central Finding: How You Represent a Signal Determines Whether You Capture It",
           "Applied Extension", 14)

    # Two-column contrast table
    section_header(s, "Same Signal — Two Representations — Opposite Results", LM, CONT_TOP)
    gs_table(s,
        ["", "hmm_state as RF Feature (continuous)", "hmm_state as Trade Gate (binary)"],
        [
            ["How used",     "Integer 0/1 passed as one of 10 features to RF tree splits",
                             "Boolean: enter only if state = 1 (stable)"],
            ["Effect on AUC","−0.122  (removing it improves AUC to 0.667)",
                             "N/A — not used in classification"],
            ["Effect on P&L","Zero — RF filter passes 28 of 29 events anyway",
                             "+0.12 Sharpe  |  +18 pp win rate"],
            ["Why?",         "RF can't recover latent structure from 56 sparse outcome labels — it memorizes the majority class",
                             "HMM imposes regime structure directly from market dynamics — no outcome labels needed"],
        ],
        LM, CONT_TOP + Inches(0.34), CW, Inches(2.50),
        col_widths=[1.50, 5.70, 5.69],
        hdr_size=10.5, body_size=11, first_col_bold=True)

    section_header(s, "Summary of All Results", LM, CONT_TOP + Inches(3.00))
    gs_table(s,
        ["Question", "Answer", "Key Number"],
        [
            ["RQ1: Does mean reversion exist?",
             "Yes — in every direction, every maturity",
             "100% sign consistency; 5Y d=0.87"],
            ["RQ2: Is it robust?",
             "Yes — 162 combinations, sign never flips",
             "73% survive FDR; range ≤ 2.83 bps"],
            ["RQ3: When is it strongest?",
             "Rising-rate regimes, 3× larger reversion",
             "Partial η²=0.28 for rate regime; vol η²=0.000"],
            ["Applied: does ML add value?",
             "HMM gate yes (+0.12 Sharpe); RF no (inoperative)",
             "Representation > model complexity"],
        ],
        LM, CONT_TOP + Inches(3.34), CW, Inches(1.60),
        col_widths=[3.0, 6.0, 3.89],
        hdr_size=10.5, body_size=11)

    rect(s, 0, CONT_BOT - Inches(0.75), W, Inches(0.75), DARK)
    textbox(s,
            "\"Given scarce labels, imposing known market structure beats trying to learn it from outcomes.\"",
            LM, CONT_BOT - Inches(0.62), CW, Inches(0.55),
            bold=True, size=14, color=WHITE, align='center')


# ─── Slide 15: Limitations & Future Work ──────────────────────────────────────
def s15_limitations():
    s = blank_slide()
    chrome(s, "Limitations Are Real — the Core Signal Remains Robust",
           "Limitations & Future Work", 15)

    lw = (CW - Inches(0.35)) / 2
    rx = LM + lw + Inches(0.35)

    section_header(s, "Current Limitations", LM, CONT_TOP)
    gs_table(s,
        ["Limitation", "Detail", "Impact"],
        [
            ["OOS window",       "27 months, 29 signal events",
                                 "Sharpe differences are directionally informative but not statistically conclusive"],
            ["ML training set",  "85 events total, 56 after label construction",
                                 "Too small for a 10-feature RF — structural failure, not tuning failure"],
            ["2Y artifacts",     "23 LSEG quoting errors in 2022",
                                 "Retained conservatively in classical tests; excluded from all ML features"],
            ["Transaction costs","Flat 1 bp round-trip assumed",
                                 "Real execution would require size-dependent market-impact modeling"],
        ],
        LM, CONT_TOP + Inches(0.34), lw, Inches(2.58),
        col_widths=[1.55, 2.85, 2.22],
        hdr_size=10, body_size=10.5)

    section_header(s, "Future Work", rx, CONT_TOP)
    gs_table(s,
        ["Extension", "Expected Value"],
        [
            ["Extend history to 2010\n(SOFR/LIBOR-OIS splice)",
             "16 years vs. 7 — enough data for conclusive strategy comparisons"],
            ["Add 1Y and 30Y maturities",
             "Tests whether regime dependence scales with maturity"],
            ["Real-time HMM filter\n(recursive, not batch)",
             "Enables live deployment without look-ahead"],
            ["Replicate in EUR and GBP\n(€STR OIS, SONIA OIS)",
             "Tests whether structural feature is global"],
            ["Pool events across maturities\nfor ML",
             "~250 events — sufficient for proper cross-validation"],
        ],
        rx, CONT_TOP + Inches(0.34), lw, Inches(2.58),
        col_widths=[2.52, 3.88],
        hdr_size=10, body_size=10.5)

    section_header(s, "What Remains Robust Despite These Limits", LM, CONT_TOP + Inches(3.10))
    bullet_block(s, [
        "Mean reversion: 100% sign consistency, 73% FDR survival, across 162 parameter combinations",
        "Regime dependence: partial η²=0.28 for rate direction — a large, stable effect",
        "HMM result: directional improvement in Sharpe (+31%) and win rate (+18 pp) with fewer trades",
    ], LM, CONT_TOP + Inches(3.44), CW, Inches(1.35), size=11.5)

    section_header(s, "Ethical Considerations", LM, CONT_TOP + Inches(4.92))
    callout(s,
            "At scale, systematic RV strategies add liquidity in stable regimes but may amplify dislocations in the stress-state windows when they trade least — monitoring systemic exposure is essential for any live deployment.",
            LM, CONT_TOP + Inches(5.26), CW, Inches(0.65), size=11, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════════════

s01_title()
s02_spread()
s03_rqs()
s04_data()
s05_rq1()
s06_event_study()
s07_rq2()
s08_rq3()
s09_ml_overview()
s10_hmm()
s11_rf()
s12_backtest()
s13_backtest_chart()
s14_finding()
s15_limitations()

OUT = "cs216_ois_presentation.pptx"
prs.save(OUT)
print(f"Saved → {OUT}  ({len(prs.slides)} slides)")
